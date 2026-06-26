from __future__ import annotations

import json
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(r"c:\wamp64\www\oase")
MANIFEST_PATH = BASE_DIR / "maquette" / "screenshots_maquette" / "manifest.json"
OUTPUT_DIR = BASE_DIR / "elaboration_rapport"
OUTPUT_PATH = OUTPUT_DIR / f"PRESENTATION_MAQUETTE_OASE_DECIDEUR_{datetime.now():%Y%m%d_%H%M%S}.pptx"

PRIMARY = RGBColor(15, 23, 42)
SECONDARY = RGBColor(30, 64, 175)
ACCENT = RGBColor(13, 148, 136)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(248, 250, 252)
BORDER = RGBColor(203, 213, 225)


@dataclass(frozen=True)
class SlideMeta:
    name: str
    title: str
    process_key: str
    process_label: str
    persona_label: str
    what_it_shows: str
    business_value: str
    pilot_message: str


PROCESS_INFO = {
    "par_processus/01_auth_et_entree": {
        "label": "1. Entree et securisation du parcours",
        "summary": "La demonstration montre comment l'acces au service reste simple pour l'usager tout en cadrant les habilitations et la continuité de compte.",
    },
    "par_processus/02_portail_beneficiaire": {
        "label": "2. Parcours du beneficiaire",
        "summary": "La maquette couvre le depot, le suivi et la consultation des avantages accordes dans un parcours compréhensible et trace.",
    },
    "par_processus/03_instruction_validation_backoffice": {
        "label": "3. Instruction et validation",
        "summary": "Les equipes de traitement disposent d'une vue de pilotage, de file de dossiers et d'ecrans de decision pour accelerer la qualite de traitement.",
    },
    "par_processus/04_workflows_regimes": {
        "label": "4. Regimes et circuits de decision",
        "summary": "Chaque regime important est presente avec son circuit propre afin de montrer que la solution s'adapte aux politiques publiques sans perdre la lisibilite du suivi.",
    },
    "par_processus/05_agences_conventions_engagements": {
        "label": "5. Agences et conventions",
        "summary": "La maquette met en scene la coordination avec les agences, les conventions et les engagements associes aux projets.",
    },
    "par_processus/06_pilotage_decideur": {
        "label": "6. Pilotage decideur",
        "summary": "Les ecrans de pilotage donnent une lecture directe des enjeux, des impacts et des arbitrages utiles pour la decision.",
    },
    "par_processus/07_tresor_rapprochements_archives": {
        "label": "7. Tresor, rapprochements et archives",
        "summary": "La maquette montre comment les validations sont rattachees aux flux financiers et a une memoire consultable des dossiers.",
    },
    "par_processus/08_institutions_specialisees": {
        "label": "8. Institutions specialisees",
        "summary": "Des vues dediees mettent en avant les besoins des institutions partenaires sans casser la coherence de l'ensemble.",
    },
    "par_processus/09_audit_controle": {
        "label": "9. Audit et controle",
        "summary": "Le dispositif de suivi permet de revisiter les decisions, de reperer les anomalies et d'organiser les missions de controle.",
    },
    "par_processus/10_open_data": {
        "label": "10. Publication et transparence",
        "summary": "La maquette prevoit une restitution publique des informations utiles, lisible et differenciee des espaces internes.",
    },
    "par_processus/11_administration_gouvernance": {
        "label": "11. Administration et gouvernance",
        "summary": "Les administrateurs disposent d'ecrans pour faire evoluer l'organisation, les regles et le pilotage sans perdre la maitrise d'ensemble.",
    },
    "par_processus/13_transverse_notifications": {
        "label": "12. Suivi transverse",
        "summary": "Le centre de notifications relie les acteurs, les dossiers prioritaires et les actions a mener.",
    },
}


SLIDE_DETAILS = {
    "01_login": ("Page d'accueil", "Met en confiance et lance la demonstration.", "Donne un point d'entree unique pour tous les publics.", "Pose le cadre d'une solution institutionnelle lisible."),
    "02_mfa": ("Verification renforcee", "Confirme l'identite avant l'acces aux dossiers.", "Protege les operations sensibles sans alourdir le discours de presentation.", "Rassure sur la maitrise des acces."),
    "03_reset-password": ("Recuperation de compte", "Permet de reprendre rapidement la main sur un compte bloque.", "Reduit les interruptions de service pour les usagers et les equipes.", "Montre une continuité de service concrete."),
    "04_activate": ("Activation de compte", "Cadre la premiere connexion et l'entree dans le service.", "Facilite le demarrage des nouveaux usagers et partenaires.", "Montre une mise en service accompagnee."),
    "06_portail-dashboard": ("Vue d'ensemble beneficiaire", "Synthese immediate de la situation du demandeur.", "Aide l'usager a comprendre ses demandes sans intermédiaire.", "Montre un service plus simple et plus lisible."),
    "07_portail-nouvelle-demande": ("Depot d'une nouvelle demande", "Guide l'usager dans la preparation d'un dossier complet.", "Ameliore la qualite des depots et limite les allers-retours.", "Montre une administration plus accessible."),
    "08_portail-demande-detail": ("Suivi detaille du dossier", "Permet de suivre l'avancement et les attentes restantes.", "Reduit l'incertitude et les relances inutiles.", "Met en avant la transparence du traitement."),
    "09_portail-exonerations-actives": ("Exonerations en cours", "Affiche clairement les avantages deja actifs.", "Donne une lecture directe des decisions produisant effet.", "Aide a mesurer le stock d'avantages accordes."),
    "10_portail-profil": ("Profil du beneficiaire", "Rassemble les informations de reference de l'usager.", "Fiabilise les echanges et les mises a jour futures.", "Montre une base commune partagee."),
    "11_backoffice-dashboard": ("Pilotage de l'instruction", "Donne la vue de charge, des priorites et des delais.", "Aide les responsables a orienter l'action quotidienne.", "Montre la capacite de pilotage operationnel."),
    "12_backoffice-dossiers": ("Portefeuille de dossiers", "Liste les dossiers a traiter avec des reperes clairs.", "Accelere le tri et la priorisation.", "Montre une meilleure maitrise des files d'attente."),
    "13_backoffice-instruction": ("Instruction du dossier", "Rassemble les pieces, controles et points d'analyse.", "Facilite un examen plus homogene des demandes.", "Soutient la qualite de la decision."),
    "14_backoffice-validation": ("Validation de la decision", "Formalise l'etape finale avant accord ou rejet.", "Cadre la responsabilite et la cohérence des arbitrages.", "Montre une prise de decision mieux securisee."),
    "15_backoffice-controle": ("Controle interne", "Met en avant les verifications et alertes utiles.", "Aide a prevenir les erreurs avant qu'elles ne se diffusent.", "Montre une administration plus fiable."),
    "16_backoffice-budget": ("Lecture budgetaire", "Relie les dossiers aux enveloppes et a leur suivi.", "Aide a piloter l'impact financier des decisions.", "Donne une vue utile pour l'arbitrage public."),
    "17_workflow-ci-otr": ("Circuit CI-OTR", "Expose un circuit de traitement adapte a ce regime.", "Montre que les regles peuvent varier selon les cas.", "Illustre la souplesse du dispositif."),
    "18_workflow-cddi": ("Circuit CDDI", "Presente un parcours specialise pour ce type de dossier.", "Evite de traiter tous les dossiers de la meme facon.", "Montre une solution adaptee aux realites de terrain."),
    "19_workflow-zone-franche": ("Circuit Zone franche", "Rend visible la specifite de ce parcours.", "Aide a expliquer les points de passage propres au regime.", "Soutient la lisibilite des politiques d'incitation."),
    "20_workflow-code-invest": ("Circuit Code des investissements", "Structure un parcours attendu par les porteurs de projet.", "Facilite la comprehension des etapes de traitement.", "Montre la capacite de soutenir l'attractivite."),
    "21_workflow-convention-miniere": ("Circuit Convention miniere", "Met en scene un traitement adapte aux enjeux miniers.", "Permet d'illustrer la prise en compte des cas complexes.", "Renforce la credibilite sectorielle de la maquette."),
    "22_agences-dashboard": ("Pilotage agence", "Donne une lecture rapide des dossiers suivis par l'agence.", "Ameliore la coordination entre promotion et instruction.", "Montre une chaine publique plus coherente."),
    "23_agences-conventions": ("Suivi des conventions", "Rassemble les conventions actives et leurs points clefs.", "Facilite la lecture des engagements pris.", "Aide au pilotage des projets strategiques."),
    "24_agences-agrements": ("Suivi des agrements", "Affiche les demandes, decisions et etats d'avancement.", "Aide a suivre un stock de dossiers heterogene.", "Montre une vision partagée entre institutions."),
    "25_agences-engagements": ("Suivi des engagements", "Suit les obligations et leur progression dans le temps.", "Aide a verifier que les promesses se traduisent en realisations.", "Met l'accent sur les resultats attendus."),
    "26_decideur-dashboard": ("Tableau de bord decideur", "Offre une synthese executive immediate.", "Permet de comprendre les volumes, tendances et tensions.", "C'est la porte d'entree naturelle pour un arbitrage."),
    "27_decideur-analyse": ("Analyse des impacts", "Met en regard les decisions et leurs effets attendus.", "Aide a objectiver les arbitrages.", "Renforce le pilotage par les resultats."),
    "28_decideur-rapport-annuel": ("Rapport annuel", "Prepare une restitution claire de l'activite et des effets.", "Facilite le dialogue avec les instances de pilotage.", "Montre une solution utile aussi pour le reporting."),
    "29_decideur-simulation": ("Simulation d'arbitrage", "Projette les effets de decisions ou scenarios.", "Aide a choisir en connaissant mieux les consequences.", "Montre une aide a la decision concrete."),
    "30_decideur-referentiel": ("Base de reference", "Rassemble les grandes regles et categories utiles a la lecture.", "Assure une interpretation commune des dossiers.", "Renforce la coherence de la decision publique."),
    "31_decideur-registre-central": ("Registre central", "Consolide la memoire des decisions et des dossiers.", "Permet une vue d'ensemble nationale.", "Soutient la gouvernance et la redevabilite."),
    "32_tresor-dashboard": ("Vue Tresor", "Relie le suivi des decisions a leur execution financiere.", "Aide a partager la meme lecture entre services.", "Montre un suivi plus complet des effets reels."),
    "33_tresor-rapprochements": ("Rapprochements", "Met en relation les dossiers et les constats financiers.", "Aide a identifier plus vite les ecarts.", "Soutient une meilleure fiabilite des comptes."),
    "34_tresor-archives": ("Archives", "Conserve une memoire accessible des dossiers traites.", "Facilite les reprises, controles et consultations.", "Montre une continuité institutionnelle."),
    "35_ministeres-dashboard": ("Vue ministeres sectoriels", "Donne une lecture adaptee aux besoins d'un ministere partenaire.", "Facilite le dialogue interinstitutionnel.", "Montre que la maquette integre toute la chaine publique."),
    "36_mae-accords-siege": ("Accords de siege", "Presente un cas specialise lie aux engagements internationaux.", "Aide a montrer que la solution couvre les situations sensibles.", "Renforce la maturite de la maquette."),
    "37_extractif-dashboard": ("Vue secteur extractif", "Adapte la lecture aux enjeux de ce secteur strategique.", "Permet de montrer une capacite de suivi sectoriel.", "Valorise la profondeur fonctionnelle de la solution."),
    "38_conedef-dashboard": ("Vue CONEDEF", "Donne une lecture propre aux besoins de concertation et suivi.", "Facilite une vision partagee entre acteurs.", "Montre la capacite d'integration institutionnelle."),
    "39_dsi-dashboard": ("Vue DSI", "Affiche une lecture de supervision et de continuité de service.", "Aide a rassurer sur l'exploitation et le suivi.", "Montre une gouvernance solide du service."),
    "40_audit-dashboard": ("Tableau de bord audit", "Donne une synthese des sujets a surveiller.", "Aide a prioriser les actions de controle.", "Renforce la confiance dans le dispositif."),
    "41_audit-journal": ("Journal de suivi", "Permet de retracer le cheminement des actions marquantes.", "Facilite les revues et investigations.", "Montre une vraie memoire des operations."),
    "42_audit-anomalies": ("Anomalies", "Rassemble les points a surveiller ou corriger.", "Permet une reaction plus rapide des responsables.", "Soutient l'amelioration continue."),
    "43_audit-dossiers": ("Dossiers a examiner", "Organise les dossiers retenus pour revue ou controle.", "Aide a structurer les campagnes de verification.", "Montre un controle mieux outille."),
    "44_audit-missions": ("Missions de controle", "Met en avant la planification et le suivi des missions.", "Facilite le pilotage des equipes de controle.", "Renforce la discipline de suivi."),
    "45_opendata-home": ("Portail public", "Ouvre une lecture grand public des informations utiles.", "Ameliore la transparence sans exposer les espaces internes.", "Montre la dimension de redevabilite."),
    "46_opendata-tableaux-de-bord": ("Tableaux de bord publics", "Restitue les grands chiffres de facon lisible.", "Aide a partager les tendances avec le public.", "Valorise l'action publique."),
    "47_opendata-datasets": ("Jeux de donnees", "Montre la mise a disposition structuree des informations ouvertes.", "Favorise la reutilisation et l'analyse externe.", "Renforce la transparence de l'Etat."),
    "48_opendata-rapports": ("Rapports publics", "Diffuse des syntheses formelles et partageables.", "Facilite la communication institutionnelle.", "Montre une ouverture maitrisee."),
    "49_admin-utilisateurs": ("Gestion des utilisateurs", "Permet d'organiser qui intervient dans la chaine.", "Aide a cadrer les responsabilites.", "Montre une gouvernance claire des acces."),
    "50_admin-roles": ("Gestion des roles", "Definit les profils et leurs droits d'action.", "Rend l'organisation plus simple a faire evoluer.", "Montre une administration durable."),
    "51_admin-connecteurs": ("Echanges avec les autres services", "Donne une vue sur les liens utiles avec l'ecosysteme public.", "Facilite la coherence entre applications.", "Montre que la solution s'insere dans l'existant."),
    "52_admin-workflow": ("Gestion des circuits", "Permet d'ajuster les etapes de traitement selon les besoins.", "Aide a accompagner les evolutions institutionnelles.", "Montre une solution adaptable dans le temps."),
    "53_admin-formulaires": ("Gestion des formulaires", "Permet de faire evoluer les formulaires sans refonte de la maquette.", "Aide a suivre les changements de politique publique.", "Montre une capacite d'adaptation rapide."),
    "54_admin-dictionnaire-o2": ("Base commune des informations", "Structure les informations de reference partagees.", "Assure une lecture homogène entre services.", "Renforce la fiabilite des analyses."),
    "55_admin-gouvernance-donnees": ("Gouvernance des donnees", "Cadre les responsabilites et les regles de qualite.", "Aide a maintenir une information fiable dans le temps.", "Montre une maitrise durable du dispositif."),
    "56_admin-requetes-dynamiques": ("Analyses a la demande", "Permet de produire des vues utiles pour les responsables.", "Accroît la reactivite face aux questions de pilotage.", "Montre une maquette utile au quotidien."),
    "57_admin-ged": ("Gestion documentaire", "Rassemble les pieces et leur consultation dans un cadre commun.", "Facilite l'examen des preuves et des justificatifs.", "Renforce la solidite des dossiers."),
    "58_admin-publication-open-data": ("Preparation des publications", "Cadre ce qui peut etre partage vers l'exterieur.", "Aide a separer information interne et information publique.", "Montre une transparence maitrisee."),
    "59_admin-regles": ("Regles de gestion", "Formalise les grands criteres qui guident les traitements.", "Aide a harmoniser les pratiques.", "Montre une gouvernance plus stable."),
    "60_admin-parametres": ("Parametrage general", "Centralise les reglages utiles a l'exploitation.", "Facilite les ajustements sans remettre en cause l'ensemble.", "Montre une solution durable et pilotable."),
    "61_admin-monitoring": ("Supervision du service", "Donne une lecture de l'etat de marche du dispositif.", "Aide a anticiper les incidents et a rassurer les responsables.", "Montre une mise en oeuvre serieuse."),
    "62_notifications": ("Centre de notifications", "Oriente chaque acteur vers ses urgences et ses prochaines actions.", "Fluidifie la coordination entre les differentes equipes.", "Montre une maquette vivante et reactive."),
}


PERSONA_LABELS = {
    "par_persona/00_entree_demo": "Entree de demonstration",
    "par_persona/P1_beneficiaire": "Beneficiaire",
    "par_persona/P2_regie_financiere": "Regie financiere",
    "par_persona/P2bis_tresor_gudef": "Tresor et GUDEF",
    "par_persona/P3_agence": "Agence de promotion",
    "par_persona/P3bis_ministeres_sectoriels": "Ministeres sectoriels",
    "par_persona/P3ter_mae_accords_siege": "MAE et accords de siege",
    "par_persona/P3quater_extractif": "Secteur extractif",
    "par_persona/P4_decideur": "Decideur",
    "par_persona/P4bis_conedef": "CONEDEF",
    "par_persona/P5_audit": "Audit et controle",
    "par_persona/P6_opendata_public": "Public",
    "par_persona/P7_admin": "Administration",
    "par_persona/P7bis_dsi_mef": "DSI",
    "par_persona/transversal": "Transversal",
}


def load_manifest() -> list[dict]:
    data = json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    captures = data["captures"]
    return [entry for entry in captures if entry["name"] != "05_mobile-mvp"]


def build_meta(entry: dict) -> SlideMeta:
    process_key = entry["process"]
    detail = SLIDE_DETAILS[entry["name"]]
    return SlideMeta(
        name=entry["name"],
        title=detail[0],
        process_key=process_key,
        process_label=PROCESS_INFO[process_key]["label"],
        persona_label=PERSONA_LABELS[entry["persona"]],
        what_it_shows=detail[1],
        business_value=detail[2],
        pilot_message=detail[3],
    )


def set_background(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, *, size=24, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = "Aptos"
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    p.alignment = align
    return box


def add_footer(slide, page_label: str) -> None:
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.0), Inches(7.05), Inches(13.33), Inches(0.42))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY
    line.line.fill.background()
    add_textbox(slide, Inches(0.35), Inches(7.09), Inches(8.5), Inches(0.22), "OASE - Maquette de presentation decideur", size=10, color=WHITE)
    add_textbox(slide, Inches(11.2), Inches(7.09), Inches(1.7), Inches(0.22), page_label, size=10, color=WHITE, align=PP_ALIGN.RIGHT)


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PRIMARY)
    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(0.8), Inches(3.4), Inches(0.55))
    banner.fill.solid()
    banner.fill.fore_color.rgb = ACCENT
    banner.line.fill.background()
    add_textbox(slide, Inches(0.92), Inches(0.92), Inches(2.9), Inches(0.2), "Presentation decideur", size=20, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.75), Inches(1.7), Inches(9.8), Inches(1.2), "Maquette OASE", size=30, bold=True, color=WHITE)
    add_textbox(
        slide,
        Inches(0.78),
        Inches(2.7),
        Inches(8.6),
        Inches(1.3),
        "Vision d'ensemble des fonctionnalites, des parcours et des ecrans clefs a presenter en COPIL ou en arbitrage.",
        size=19,
        color=RGBColor(226, 232, 240),
    )
    add_textbox(
        slide,
        Inches(0.78),
        Inches(5.45),
        Inches(5.4),
        Inches(0.8),
        f"Edition du {datetime.now():%d/%m/%Y}",
        size=16,
        color=RGBColor(191, 219, 254),
    )


def add_summary_slide(prs: Presentation, metas: list[SlideMeta]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHT)
    add_textbox(slide, Inches(0.65), Inches(0.45), Inches(6), Inches(0.5), "Sommaire", size=28, bold=True, color=PRIMARY)
    add_textbox(
        slide,
        Inches(0.65),
        Inches(0.95),
        Inches(10.8),
        Inches(0.45),
        "Lecture simple pour decideur: chaque ecran est introduit par une slide qui explique ce qu'il faut retenir.",
        size=15,
        color=MUTED,
    )
    labels = []
    seen = set()
    for meta in metas:
        if meta.process_key not in seen:
            seen.add(meta.process_key)
            labels.append(PROCESS_INFO[meta.process_key]["label"])
    top = 1.55
    for index, label in enumerate(labels, start=1):
        pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(top), Inches(0.6), Inches(0.36))
        pill.fill.solid()
        pill.fill.fore_color.rgb = SECONDARY
        pill.line.fill.background()
        add_textbox(slide, Inches(0.88), Inches(top + 0.05), Inches(0.25), Inches(0.2), str(index), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.45), Inches(top - 0.02), Inches(10.7), Inches(0.35), label, size=18, bold=True, color=TEXT)
        top += 0.47
    add_textbox(
        slide,
        Inches(0.7),
        Inches(6.25),
        Inches(11.8),
        Inches(0.45),
        f"Volume du livrable: {len(metas)} ecrans couverts hors mobile, chacun avec une slide d'explication et une slide de capture.",
        size=14,
        color=MUTED,
    )


def add_section_slide(prs: Presentation, process_key: str) -> None:
    info = PROCESS_INFO[process_key]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, SECONDARY)
    add_textbox(slide, Inches(0.72), Inches(1.25), Inches(10.6), Inches(0.7), info["label"], size=28, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.75), Inches(2.25), Inches(10.9), Inches(1.5), info["summary"], size=20, color=RGBColor(226, 232, 240))
    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(5.5), Inches(2.7), Inches(0.5))
    tag.fill.solid()
    tag.fill.fore_color.rgb = ACCENT
    tag.line.fill.background()
    add_textbox(slide, Inches(1.02), Inches(5.63), Inches(2.2), Inches(0.2), "Sequence de presentation", size=14, bold=True, color=WHITE)


def add_feature_slide(prs: Presentation, meta: SlideMeta) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHT)
    add_textbox(slide, Inches(0.65), Inches(0.42), Inches(9.5), Inches(0.6), meta.title, size=27, bold=True, color=PRIMARY)
    add_textbox(slide, Inches(0.67), Inches(0.98), Inches(5), Inches(0.3), meta.process_label, size=13, color=SECONDARY)
    add_textbox(slide, Inches(5.3), Inches(0.98), Inches(2.4), Inches(0.3), f"Public: {meta.persona_label}", size=13, color=ACCENT, align=PP_ALIGN.RIGHT)

    card1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(1.55), Inches(3.85), Inches(3.95))
    card1.fill.solid()
    card1.fill.fore_color.rgb = WHITE
    card1.line.color.rgb = BORDER
    add_textbox(slide, Inches(1.0), Inches(1.85), Inches(3.1), Inches(0.35), "Ce que montre l'ecran", size=18, bold=True, color=PRIMARY)
    add_textbox(slide, Inches(1.0), Inches(2.35), Inches(3.05), Inches(2.7), meta.what_it_shows, size=20, color=TEXT)

    card2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.55), Inches(3.85), Inches(3.95))
    card2.fill.solid()
    card2.fill.fore_color.rgb = WHITE
    card2.line.color.rgb = BORDER
    add_textbox(slide, Inches(5.08), Inches(1.85), Inches(3.1), Inches(0.35), "Valeur pour l'action publique", size=18, bold=True, color=PRIMARY)
    add_textbox(slide, Inches(5.08), Inches(2.35), Inches(3.05), Inches(2.7), meta.business_value, size=20, color=TEXT)

    card3 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.88), Inches(1.55), Inches(3.7), Inches(3.95))
    card3.fill.solid()
    card3.fill.fore_color.rgb = WHITE
    card3.line.color.rgb = BORDER
    add_textbox(slide, Inches(9.15), Inches(1.85), Inches(2.95), Inches(0.35), "Message a retenir", size=18, bold=True, color=PRIMARY)
    add_textbox(slide, Inches(9.15), Inches(2.35), Inches(2.9), Inches(2.7), meta.pilot_message, size=20, color=TEXT)

    ribbon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(6.0), Inches(11.8), Inches(0.58))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = ACCENT
    ribbon.line.fill.background()
    add_textbox(slide, Inches(1.0), Inches(6.16), Inches(11.2), Inches(0.2), "La slide suivante presente la capture exacte de l'ecran correspondant.", size=14, bold=True, color=WHITE)


def add_capture_slide(prs: Presentation, meta: SlideMeta, image_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_textbox(slide, Inches(0.55), Inches(0.28), Inches(10.4), Inches(0.45), meta.title, size=24, bold=True, color=PRIMARY)
    add_textbox(slide, Inches(0.58), Inches(0.72), Inches(8.8), Inches(0.25), f"{meta.process_label} - {meta.persona_label}", size=12, color=MUTED)

    frame_left = Inches(0.48)
    frame_top = Inches(1.0)
    frame_width = Inches(12.35)
    frame_height = Inches(5.72)
    padding = Inches(0.14)

    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, frame_left, frame_top, frame_width, frame_height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = LIGHT
    frame.line.color.rgb = BORDER

    available_width = frame_width - (padding * 2)
    available_height = frame_height - (padding * 2)

    with Image.open(image_path) as image:
        image_width, image_height = image.size

    image_ratio = image_width / image_height
    available_ratio = available_width / available_height

    if image_ratio > available_ratio:
        picture_width = available_width
        picture_height = available_width / image_ratio
    else:
        picture_height = available_height
        picture_width = available_height * image_ratio

    picture_left = frame_left + ((frame_width - picture_width) / 2)
    picture_top = frame_top + ((frame_height - picture_height) / 2)

    slide.shapes.add_picture(str(image_path), picture_left, picture_top, width=picture_width, height=picture_height)


def add_page_numbers(prs: Presentation) -> None:
    total = len(prs.slides)
    for index, slide in enumerate(prs.slides, start=1):
        add_footer(slide, f"Page {index} / {total}")


def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    captures = load_manifest()
    metas = [build_meta(entry) for entry in captures]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_summary_slide(prs, metas)

    current_process = None
    for entry, meta in zip(captures, metas):
        if meta.process_key != current_process:
            current_process = meta.process_key
            add_section_slide(prs, current_process)
        add_feature_slide(prs, meta)
        add_capture_slide(prs, meta, Path(entry["files"][1]))

    add_page_numbers(prs)
    prs.save(OUTPUT_PATH)
    print(OUTPUT_PATH)


if __name__ == "__main__":
    main()
